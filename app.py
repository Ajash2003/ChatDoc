import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain.vectorstores import FAISS
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
from fpdf2 import FPDF2 as FPDF

load_dotenv()

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_ppt_text(ppt_docs):
    text = ""
    for ppt in ppt_docs:
        presentation = Presentation(ppt)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
    return text

def get_doc_text(doc_docs):
    text = ""
    for doc in doc_docs:
        doc_file = docx.Document(doc)
        for para in doc_file.paragraphs:
            text += para.text + " "
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=10000,
        chunk_overlap=1000
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, 
    if the answer is given in points make sure to give the points in different lines, 
    if the answer is not in the provided context just say, "Answer cannot be found", don't provide the wrong answer.
    Context: \n{context}\n
    Question: \n{question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()

    response = chain(
        {"input_documents": docs, "question": user_question},
        return_only_outputs=True
    )
    
    answer = response.get("output_text", "").strip()
    if not answer:
        return "Answer cannot be found"
    return answer

def generate_pdf(qa_pairs):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    
    # Use the built-in Arial font
    pdf.add_font('Arial', '', uni=True)
    pdf.set_font('Arial', size=12)

    # Title
    pdf.cell(200, 10, txt="Chat Conversation", ln=True, align='C')
    pdf.ln(10)

    for idx, (question, answer) in enumerate(qa_pairs, start=1):
        try:
            # Question
            pdf.set_font('Arial', 'B', size=12)  # Bold for questions
            pdf.cell(0, 10, txt=f"Q{idx}:", ln=True)
            pdf.set_font('Arial', size=12)  # Regular for content
            pdf.multi_cell(0, 10, txt=question)
            pdf.ln(5)

            # Answer
            pdf.set_font('Arial', 'B', size=12)  # Bold for answer label
            pdf.cell(0, 10, txt=f"A{idx}:", ln=True)
            pdf.set_font('Arial', size=12)  # Regular for content
            pdf.multi_cell(0, 10, txt=answer)
            pdf.ln(10)

        except Exception as e:
            # If there's still an encoding error, fall back to ASCII
            question = ''.join(char for char in question if ord(char) < 128)
            answer = ''.join(char for char in answer if ord(char) < 128)
            
            pdf.set_font('Arial', 'B', size=12)
            pdf.cell(0, 10, txt=f"Q{idx}:", ln=True)
            pdf.set_font('Arial', size=12)
            pdf.multi_cell(0, 10, txt=question)
            pdf.ln(5)
            
            pdf.set_font('Arial', 'B', size=12)
            pdf.cell(0, 10, txt=f"A{idx}:", ln=True)
            pdf.set_font('Arial', size=12)
            pdf.multi_cell(0, 10, txt=answer)
            pdf.ln(10)

    file_path = "chat_conversation.pdf"
    pdf.output(file_path)
    return file_path

def main():
    st.set_page_config(page_title="ChatDoc", page_icon=":books:")
    st.markdown("""
    <style>
    .main {
        padding: 20px;
    }
    .stTextInput > div > div > input {
        border: 1px solid #ddd;
        border-radius: 10px;
        padding: 10px;
        font-size: 18px;
    }
    .stButton button {
        background-color: #4CAF50;
        color: white;
        padding: 10px 20px;
        border-radius: 10px;
        border: none;
        font-size: 18px;
        cursor: pointer;
    }
    .stButton button:hover {
        background-color: #45a049;
    }
    .question-box {
        border: 2px solid #2196F3;
        border-radius: 10px;
        padding: 15px;
        margin-bottom: 20px;
        background-color: #2b313e;
        box-shadow: 0 4px 8px rgba(0,0,0,0.1);
    }
    .answer-box {
        border: 2px solid #4CAF50;
        border-radius: 10px;
        padding: 20px;
        margin-bottom: 20px;
        background-color: #475063;
        box-shadow: 0 4px 8px rgba(0,0,0,0.1);
    }
    .sidebar .sidebar-content {
        background-color: #2b313e;
        padding: 20px;
        border-radius: 10px;
    }
    .stFileUploader {
        background-color: #2b313e;
        border: 1px solid #2b313e;
        border-radius: 10px;
    }
    .header {
        font-size: 24px;
        font-weight: bold;
        padding: 20px 0;
    }
    .footer {
        font-size: 14px;
        color: #888;
        padding: 10px 0;
        text-align: center;
    }
    </style>
    """, unsafe_allow_html=True)

    st.header("ChatDoc :books:")
    st.markdown('<div class="header">Chat with your Documents Here</div>', unsafe_allow_html=True)

    # Initialize session state for storing questions and answers
    if "qa_pairs" not in st.session_state:
        st.session_state.qa_pairs = []

    if "processed" not in st.session_state:
        st.session_state.processed = False

    with st.sidebar:
        st.subheader("Your Uploaded Files : ")
        doc_files = st.file_uploader("Upload your PDFs, PPTs, or DOCs here:", accept_multiple_files=True, type=['pdf', 'pptx', 'docx'])
        if st.button('Process'):
            if doc_files:
                with st.spinner('Processing...'):
                    raw_text = ""
                    pdf_files = [file for file in doc_files if file.name.endswith('.pdf')]
                    ppt_files = [file for file in doc_files if file.name.endswith('.pptx')]
                    docx_files = [file for file in doc_files if file.name.endswith('.docx')]
                    
                    if pdf_files:
                        raw_text += get_pdf_text(pdf_files)
                    if ppt_files:
                        raw_text += get_ppt_text(ppt_files)
                    if docx_files:
                        raw_text += get_doc_text(docx_files)

                    text_chunks = get_text_chunks(raw_text)
                    get_vector_store(text_chunks)
                    st.session_state.processed = True
                    st.success("Processing complete. You can now ask questions.")
            else:
                st.warning("Please upload at least one PDF, PPT, or DOC document before processing.")

    if st.session_state.processed:
        col1, col2 = st.columns([8, 2])
        with col1:
            user_question = st.text_input("Ask a question:", "", key="question", help="Type your question here")
        with col2:
            if st.button("Clear", key="clear_chat"):
                st.session_state.qa_pairs = []

        if user_question:
            with st.spinner('Fetching answer...'):
                answer = user_input(user_question)
                if not st.session_state.qa_pairs or st.session_state.qa_pairs[-1][0] != user_question:
                    st.session_state.qa_pairs.append((user_question, answer))

        # Display previous questions and answers
        for question, answer in reversed(st.session_state.qa_pairs):
            st.markdown(f'<div class="question-box"><strong>Question:</strong><br>{question}</div>', unsafe_allow_html=True)
            st.markdown(f'<div class="answer-box"><strong>Answer:</strong><br>{answer}</div>', unsafe_allow_html=True)

        if st.session_state.qa_pairs:
            if st.button("Download Conversation as PDF"):
                pdf_path = generate_pdf(st.session_state.qa_pairs)
                with open(pdf_path, "rb") as file:
                    st.download_button("Download PDF", file, file_name="chat_conversation.pdf")
    else:
        st.info("Please upload and process a document to start asking questions.")

if __name__ == '__main__':
    main()
